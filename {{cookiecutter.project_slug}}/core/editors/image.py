"""

Image editors, layout spread generators, and image group transformers

"""

import os
import io
import copy
import imageio
from collections import Counter
import requests
import numpy as np
from pptx import Presentation
from pptx.util import Inches
from PIL import Image as PilImage
from PIL import ImageChops
from config import Config
from core.editors.color import Color

PilImage.MAX_IMAGE_PIXELS = 933120000


def download_image(image_url, target_fname=None):
    """
    Download an image file from the url and save it to a specified location.
    If the target location is unspecified, then save it to a temp dir
    :param image_url:
    :type image_url:
    :param target_fname:
    :type target_fname:
    :return:
    :rtype:
    """
    if target_fname is None:
        target_fname = os.path.join(Config.TEMP_DIR, image_url.split('/')[-1].split('?')[0])
    img_data = requests.get(image_url).content
    with open(target_fname, 'wb') as fp:
        fp.write(img_data)
    return target_fname


class Image:
    def __init__(self, img: PilImage.Image = None, img_fpath=None, convert_transparency=True):
        """
        Set up an Image as an extended version of a PIL image using either a PIL image or image path.
        :param img: Target image to be transformed
        :type img: PIL.Image
        :param img_fpath: Path to image to be loaded and transformed
        :type img_fpath: str
        :param convert_transparency: (Optional) Convert PNG images to RGBA. Default True.
        :type convert_transparency: boolean
        """
        self.img_fpath = img_fpath
        if img is None and img_fpath is None:
            raise ValueError("either img or img_path or both must be specified.")
        if isinstance(img, PilImage.Image):
            self.img = img
        elif isinstance(img, Image):
            self.img = img.img
            self.img_fpath = img.img_fpath
        elif img is None and os.path.isfile(img_fpath):
            self.img = self.open_as_pil_image(img_fpath=img_fpath, convert_transparency=convert_transparency)
        else:
            raise ValueError(f"either a valid img or img_fpath must be specified")
        # Temporary instance used for saving and restoring data around transformations
        self.backup_instance = None

    def copy(self):
        """
        Return a deep copy of the image
        :return:
        :rtype:
        """
        return Image(img=copy.deepcopy(self.img), img_fpath=self.img_fpath)

    def backup(self):
        """
        Return a deep copy of the image
        :return:
        :rtype:
        """
        self.backup_instance = self.copy()

    def restore(self):
        """
        Return a deep copy of the image
        :return:
        :rtype:
        """
        backup_copy = self.backup_instance.copy()
        self.img = backup_copy.img
        self.img_fpath = backup_copy.img_fpath

    @staticmethod
    def open_as_pil_image(img_fpath, convert_transparency=True):
        """
        Open an image file and load it as a PIL image
        :param img_fpath: The path to the image file
        :param convert_transparency: (Optional) Convert PNG images to RGBA. Default True.
        """
        with open(img_fpath, 'rb') as fp:
            img = PilImage.open(fp)
            img.load()
            if convert_transparency is True and img_fpath.lower().endswith('png'):
                img = img.convert("RGBA")
        return img

    @staticmethod
    def image_extensions():
        """
        Get a list of valid extensions for images based on the PIL library compatibility
        :return: List of valid image extensions
        :rtype: list of str
        """
        return ['bmp', 'eps', 'gif', 'ico', 'im', 'jpg', 'jpeg', 'j2k', 'png', 'tiff', 'tif', 'webp']

    @property
    def size(self):
        return self.img.size

    @property
    def width(self):
        return self.size[0]

    @property
    def height(self):
        return self.size[1]

    @property
    def format(self):
        return self.img.mode

    @property
    def is_transparent(self):
        if self.img.mode == "P":
            transparent = self.img.info.get("transparency", -1)
            for _, index in self.img.getcolors():
                if index == transparent:
                    return True
        elif self.img.mode == "RGBA":
            extrema = self.img.getextrema()
            if extrema[3][0] < 255:
                return True
        return False

    @property
    def bg_color(self):
        """
        Get a tuple of 0-255 values for the image background color by sampling the most common color
        :return: tuple of 0-255 values for (R, G, B, Alpha)
        :rtype: tuple
        """
        if self.is_transparent:
            return 0, 0, 0, 0
        palette = self.color_palette
        if len(palette) > 0:
            bg_color = Color.parse_hex(self.color_palette[0])
        else:
            bg_color = (255, 255, 255, 255)
        if len(bg_color) == 3:
            bg_color = (bg_color[0], bg_color[1], bg_color[2], 0)
        return bg_color

    @property
    def color_palette(self):
        img_palette = Counter(self.img.getdata())
        # Remove transparent pixels if an alpha channel is included
        img_palette = {color: freq for color, freq in img_palette.items() if
                       (len(color) > 3 and color[3] > 50) or (len(color) == 3)}
        # Remove black and white pixels
        img_palette = {color: freq for color, freq in img_palette.items() if color[:3] not in ((0, 0, 0))}
        # Format the colors to hex and sort by frequency
        img_palette = {"#{0:02x}{1:02x}{2:02x}".format(*color[:3]).upper(): freq for color, freq in
                       img_palette.items()}
        img_palette = dict(sorted(img_palette.items(), key=lambda item: -item[1]))
        return list(img_palette.keys())

    @staticmethod
    def similarity(img_1: PilImage.Image, img_2: PilImage.Image):
        """ Get the similarity between two images """
        if img_2.size != img_1.size:
            img_2 = copy.deepcopy(img_2)
            img_2 = img_2.resize(img_1.size)
            img_2 = img_2.convert(img_1.mode)
        try:
            return 1 - (np.array(ImageChops.difference(img_1, img_2).getdata()).mean() / 255)
        except ValueError:
            # ImageChops will throw a ValueError if the images don't match, so return a similarity of 0.0
            return 0.0

    @classmethod
    def file_similarity(cls, img_1: str, img_2: str):
        """ Get the similarity between two images """
        return cls.similarity(cls.open_as_pil_image(img_1), cls.open_as_pil_image(img_2))

    def resize_to_width(self, width):
        """ Resize an image to a specific width while preserving aspect ratio"""
        width_percent = (width / float(self.img.size[0]))
        height_percent = int((float(self.img.size[1]) * float(width_percent)))
        self.img = self.img.resize((width, height_percent), PilImage.LANCZOS)
        return self

    def clamp_to_width(self, max_width):
        """ Clamp images that exceed a max width to that width"""
        if self.img.size[0] > max_width:
            return self.resize_to_width(max_width)
        else:
            return self

    def resize_to_height(self, height):
        """ Resize an image to a specific height while preserving aspect ratio"""
        height_percent = (height / float(self.img.size[1]))
        width_percent = int((float(self.img.size[0]) * float(height_percent)))
        self.img = self.img.resize((width_percent, height), PilImage.LANCZOS)
        return self

    def rotate(self, angle, expand=False, center=None, translate=None, fillcolor=None):
        self.img = self.img.rotate(angle=angle, expand=expand, center=center, translate=translate, fillcolor=fillcolor)
        return self

    def crop_and_resize(self, crop=None, width=None, height=None):
        """
        Crop the image if specified with a crop tuple of (left, upper, right, lower) and resize to a target width
        :param crop: Cropping coordinates as a tuple of (left, upper, right, lower)
        :type crop: tuple of int
        :param width: Final width to resize to after cropping
        :type width: int
        :param height: Final height to resize to after cropping
        :type height: int
        :return: Cropped and resized image
        :rtype: PIL image
        """
        if isinstance(crop, (list, tuple)) and len(crop) == 4:
            self.img = self.img.crop(crop)
        if isinstance(width, int):
            self.resize_to_width(width=width)
        if isinstance(height, int):
            self.resize_to_height(height=height)
        return self

    def expand_bg(self, target_size, bg_color=None):
        """
        Expand an image, filling the background with the specified color
        """
        bg_color = bg_color or self.bg_color
        new_im = PilImage.new("RGB", target_size, bg_color[:3])
        new_im.paste(self.img, ((target_size[0] - self.width) // 2,
                                (target_size[1] - self.height) // 2))
        self.img = new_im
        return self

    def bg_to_transparent(self, bg_color=None):
        """
        Convert a PIL image with a specified or white background into an image with a transparent background.
        :param bg_color: (Optional) Tuple of RGB colors for image background. Defaults to (255, 255, 255)
        :type bg_color: tuple of uint8
        :return: new image with transparent background
        :rtype: PIL.Image
        """
        bg_color = bg_color or self.bg_color
        self.img = self.img.convert("RGBA")
        img_data = self.img.getdata()
        new_data = []
        for item in img_data:
            if item[0] == bg_color[0] and item[1] == bg_color[1] and item[2] == bg_color[2]:
                new_data.append(bg_color)
            else:
                new_data.append(item)
        self.img.putdata(new_data)
        return self

    def transparency_to_bg_color(self, bg_color=None):
        """
        Only process if image has transparency (https://stackoverflow.com/a/1963146)
        :param bg_color: (Optional) Tuple of RGB colors for image background. Defaults to (255, 255, 255)
        :type bg_color: tuple of uint8
        :return: PIL image with transparency removed
        :rtype: PIL.Image
        """
        bg_color = bg_color or self.bg_color
        if bg_color == (0, 0, 0, 0):
            bg_color = (255, 255, 255, 255)
        if self.img.mode in ('RGBA', 'LA') or (self.img.mode == 'P' and 'transparency' in self.img.info):
            # Need to convert to RGBA if LA format due to a bug in PIL (http://stackoverflow.com/a/1963146)
            alpha = self.img.convert('RGBA').split()[-1]
            # Create a new background image of our matt color.
            # Must be RGBA because paste requires both images have the same format
            # (http://stackoverflow.com/a/8720632  and  http://stackoverflow.com/a/9459208)
            bg_color = (bg_color[0], bg_color[1], bg_color[2], 255)
            bg = PilImage.new("RGB", self.img.size, bg_color)
            bg.paste(self.img, mask=alpha)
            background = PilImage.new("RGB", bg.size, (255, 255, 255))
            background.paste(bg, mask=alpha)  # 3 is the alpha channel
            self.img = background
            return self
        else:
            return self

    def paste(self, overlay: PilImage.Image, paste_coordinates, page_width=690):
        """
        Paste an overlay onto this image
        :param overlay:
        :type overlay: PilImage.Image
        :param paste_coordinates:
        :type paste_coordinates:
        :param page_width:
        :type page_width:
        :return:
        :rtype:
        """
        new_overlay = copy.deepcopy(overlay)
        new_overlay = new_overlay.convert("RGBA")
        self.img = self.img.convert("RGBA")
        new_overlay = Image(img=new_overlay).resize_to_width(page_width).img
        self.img.paste(new_overlay, (paste_coordinates[0], paste_coordinates[1]), new_overlay)
        return self

    def paste_to_bg(self, background: PilImage.Image, paste_coordinates=None, page_width=690):
        """
        TODO - Deprecated - paste this image onto a background and return the background
        :param background:
        :type background: PilImage.Image
        :param paste_coordinates:
        :type paste_coordinates:
        :param page_width:
        :type page_width:
        :return:
        :rtype:
        """
        pc = paste_coordinates or [80, 100, 820, 100]
        new_background = copy.deepcopy(background)
        new_background.convert("RGB")
        self.resize_to_width(page_width)
        new_background.paste(self.img, (pc[0], pc[1]), self.img.convert("RGBA"))
        self.img = new_background
        return self

    def convert(self, convert_str):
        """
        Convert the image into a different format, like RGB or RGBA
        :param convert_str:
        :type convert_str:
        :return:
        :rtype:
        """
        self.img = self.img.convert(convert_str)
        return self

    def get_fname(self, base_name=None, out_dir=None, width=None, is_transparent=None, new_ext=None):
        """
        Get
        :param base_name:
        :type base_name:
        :param out_dir:
        :type out_dir:
        :param width:
        :type width:
        :param is_transparent:
        :type is_transparent:
        :param new_ext: (Optional) Specify an extension to save the image as
        :type new_ext:
        :return:
        :rtype:
        """
        if width is not None:
            width = f"--{width}"
        else:
            width = ''
        if is_transparent is None:
            tra = ''
        else:
            tra = '_t' if is_transparent is True else ''
        if base_name is None and self.img_fpath is not None:
            base_name = self.img_fpath
        if out_dir is not None:
            base_name = os.path.join(out_dir, base_name)
        ext = new_ext or os.path.splitext(base_name)[1]
        if not ext.startswith('.'):
            ext = f'.{ext}'
        return f"{os.path.splitext(base_name)[0]}{width}{tra}{ext}"

    def save(self, out_path, optimize=True, quality=100):
        """
        Save the PIL image to a target output path
        :param out_path:
        :type out_path:
        :param optimize: Take a second pass over the image to reduce size
        :type optimize: float
        :param quality: Specify save quality for lossy formats like JPEG
        :type quality: float
        :return:
        :rtype:
        """
        if any(out_path.lower().endswith(ext) for ext in ('jpg', 'jpeg', 'webp')) and self.is_transparent:
            self.transparency_to_bg_color()
        print(f'image saved to {out_path}')
        return self.img.save(out_path, optimize=optimize, quality=quality)

    def compress(self, base_name=None, out_dir=None, max_width=1920, save_webp=False):
        """
        Save the image
        :param base_name:
        :type base_name:
        :param out_dir:
        :type out_dir:
        :param max_width:
        :type max_width:
        :return:
        :rtype:
        """
        ext = os.path.splitext(self.img_fpath)[-1]
        self.backup()
        if self.width > max_width:
            self.resize_to_width(max_width)
        print('new image', self.size, self.img.mode, self.img_fpath)
        self.save(out_path=self.get_fname(base_name=base_name, out_dir=out_dir, new_ext=ext), quality=30)
        if save_webp is True:
            self.save(out_path=self.get_fname(base_name=base_name, out_dir=out_dir, new_ext='.webp'))
        # Reload the original state before the width transform
        self.restore()
        return self

    def process(self, out_fname, **kwargs):
        # Apply transformations on the base image
        if 'crop' in kwargs.keys():
            pass
        if 'background_image' in kwargs.keys():
            pass
        if 'is_transparent' in kwargs.keys():
            pass
        if 'expand_bg_scale' in kwargs.keys():
            pass
        # Rescale after transformations
        if 'width' in kwargs.keys():
            pass
        elif 'height' in kwargs.keys():
            pass
        elif 'scale' in kwargs.keys():
            pass
        # Save, compress, and convert as specified
        if 'save_path' in kwargs.keys():
            pass
        if 'do_compress' in kwargs.keys():
            pass
        if 'do_primitive' in kwargs.keys():
            pass
        return self


class ImageGroup:
    """
    Manage an ordered list of PIL images
    """

    def __init__(self, image_list=None, out_dir=None):
        """
        :param image_list: List of images as strings or PIL images
        :type image_list: list of str or list of PIL image
        :param out_dir: Output directory where results will be stored
        :type out_dir: str
        """
        # Open the images as PIL images if passed as a list of str
        if isinstance(image_list[0], str):
            self.image_list = [Image(img_fpath=img) for img in image_list]
        elif isinstance(image_list[0], Image):
            self.image_list = image_list
        elif isinstance(image_list[0], PilImage.Image):
            self.image_list = [Image(img=img) for img in image_list]
        else:
            raise ValueError('image_list must be specified as either a list of str fpaths or PIL images')
        # Set up out_dir and interim_dir
        if out_dir is None or not os.path.isdir(out_dir):
            out_dir = Config.TEST_GALLERY_DIR
        interim_dir = Config.TEST_GALLERY_DIR
        self.out_dir = out_dir
        self.interim_dir = interim_dir

    @staticmethod
    def filter_image_list(image_list):
        """
        Reduce a list of paths to only the ones with valid image extensions
        :param image_list: List of image fnames or fpaths
        :type image_list: list of str
        :return: List of image fnames filtered by extension
        :rtype: list of str
        """
        return [img for img in image_list
                if any([img.lower().endswith(img_ext) for img_ext in Image.image_extensions()])]

    def apply(self, in_func, *args, **kwargs):
        """
        Apply a function to the list of images.
        :param in_func:
        :type in_func:
        :return:
        :rtype:
        """
        self.image_list = [in_func(img, *args, **kwargs) for img in self.image_list]

    def to_gif(self, out_fpath, fps=8, width=300, crop=None):
        """
        Turn the list of images into an animated gif
        :param out_fpath: Output file path for the GIF image
        :param fps: Frames-per-second for the gif animation
        :param width: Final width to resize to after cropping
        :param crop: Cropping coordinates as a tuple of (x0, y0, x1, y1)
        :return: None
        :rtype: None
        """
        with imageio.get_writer(out_fpath, mode='I', fps=fps, subrectangles=True) as writer:
            for image in self.image_list:
                # Convert to PIL
                image = image.crop_and_resize(crop=crop, width=width)
                # Convert from PIL
                image = np.asarray(image.img)
                writer.append_data(image)
        writer.close()
        return out_fpath

    def to_slideshow(self, out_path=None, message_list=None, in_dir='data/toolkit/image_to_ppt/in'):
        """
        Create a powerpoint presentation with a list of images.
        :param out_path:
        :type out_path:
        :param message_list: List of text to write on each slide. Used to id filenames in the doc
        :type message_list: list
        :param in_dir:
        :type in_dir:
        :return:
        :rtype:
        """
        if out_path is None:
            out_path = 'data/toolkit/image_to_ppt/out/out_ppt.ppt'
        if message_list is None:
            message_list = ['' for img in self.image_list]
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        for idx, image in enumerate(self.image_list):
            slide = prs.slides.add_slide(blank_slide_layout)
            with io.BytesIO() as output:
                image = image.transparency_to_bg_color()
                image.img.save(output, format="GIF")
                hpercent = (5.5 / float(image.size[1]))
                width_size = float(image.size[0]) * float(hpercent)
                if width_size > 9.0:
                    width_percent = (9.0 / float(image.size[0]))
                    height_size = int((float(image.size[1]) * float(width_percent)))
                    slide.shapes.add_picture(output, Inches(0.5), Inches(1),
                                             width=Inches(9.0), height=Inches(height_size))
                else:
                    slide.shapes.add_picture(output, Inches((9.0 - width_size) / 2), Inches(1),
                                             width=Inches(width_size), height=Inches(5.5))
                if message_list is not None and len(message_list) == len(self.image_list):
                    txt_box = slide.shapes.add_textbox(left=Inches(1), top=Inches(0.3),
                                                       width=Inches(5), height=Inches(1))
                    tf = txt_box.text_frame
                    tf.text = message_list[idx]
        prs.save(out_path)
        return out_path

    def to_pdf(self, out_fname):
        """
        Save the batch of images to a PDF format using PIL
        :param out_fname:
        :type out_fname:
        :return:
        :rtype:
        """
        if self.image_list[0].is_transparent is True:
            self.image_list = [img.transparency_to_bg_color((255, 255, 255)) for img in self.image_list]
        self.image_list[0].img.convert('RGB').save(out_fname, "PDF", resolution=339, save_all=True,
                                                   append_images=[each_image.img.convert('RGB')
                                                                  for each_image in self.image_list[1:]])
        return out_fname

    def process(self, out_fnames=None, **kwargs):
        self.image_list = [img.process(out_fname=out_fnames[idx], **kwargs)
                           for idx, img in enumerate(self.image_list)]
        return self


def get_image_group(image_list=None, image_dir=None, out_dir=None):
    """
    Generate an ImageGroup from a list of images or target directory.
    :param image_list:
    :type image_list:
    :param image_dir:
    :type image_dir:
    :param out_dir:
    :type out_dir:
    :return:
    :rtype:
    """
    if image_list is None and image_dir is None:
        raise ValueError("either image_list or image_dir or both must be specified.")
    if isinstance(image_list, (list, tuple)) and image_dir is None:
        return ImageGroup(image_list=image_list, out_dir=out_dir)
    elif image_list is None and os.path.isdir(image_dir):
        image_list = [os.path.join(image_dir, img) for img in os.listdir(image_dir)]
        image_list = ImageGroup.filter_image_list(image_list)
    elif isinstance(image_list, (list, tuple)) and os.path.isdir(image_dir):
        if not os.path.isfile(image_list[0]):
            image_list = [os.path.join(image_dir, img) for img in image_list]
            image_list = ImageGroup.filter_image_list(image_list)
    else:
        raise ValueError(f"ImageGroup cannot be built: os.path.isdir(image_dir)={os.path.isdir(image_dir)} "
                         f"isinstance(image_list, (list, tuple))={isinstance(image_list, (list, tuple))}")
    return ImageGroup(image_list=image_list, out_dir=out_dir)
